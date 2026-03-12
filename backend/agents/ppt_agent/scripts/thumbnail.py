"""Generate thumbnail grid of slides for template analysis.

Creates a composite image showing all slides from a PPTX with filenames as labels.
Useful for analyzing template layouts before editing.

Usage:
    python thumbnail.py <pptx_file> [output_prefix] [--cols N]

Examples:
    python thumbnail.py template.pptx                 # → thumbnails.jpg
    python thumbnail.py template.pptx my_output       # → my_output.jpg
    python thumbnail.py template.pptx --cols 4         # 4 columns
"""

import argparse
import os
import re
import sys
import tempfile
from pathlib import Path

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    print("Error: Pillow is required. Install with: pip install Pillow")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    Presentation = None


def create_thumbnails(
    pptx_path: str,
    output_prefix: str = "thumbnails",
    cols: int = 3,
    max_slides: int = 12,
    thumb_width: int = 400,
) -> str:
    """Create a thumbnail grid of slides from a PPTX.

    Args:
        pptx_path: Path to the PPTX file
        output_prefix: Output filename prefix (default: thumbnails)
        cols: Number of columns in the grid
        max_slides: Maximum number of slides to include
        thumb_width: Width of each thumbnail in pixels

    Returns:
        Path to the output image
    """
    if Presentation is None:
        return _create_thumbnails_fallback(pptx_path, output_prefix, cols, max_slides, thumb_width)

    # Read presentation
    prs = Presentation(pptx_path)
    slides = list(prs.slides)[:max_slides]

    if not slides:
        print("No slides found")
        return None

    # Calculate dimensions
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    aspect_ratio = slide_h / slide_w if slide_w > 0 else 0.75
    thumb_height = int(thumb_width * aspect_ratio)

    label_height = 25
    padding = 10
    rows = (len(slides) + cols - 1) // cols

    grid_width = cols * (thumb_width + padding) + padding
    grid_height = rows * (thumb_height + label_height + padding) + padding

    # Create grid image
    grid = Image.new("RGB", (grid_width, grid_height), "white")
    draw = ImageDraw.Draw(grid)

    # Try to load a font
    try:
        font = ImageFont.truetype("arial.ttf", 14)
    except (IOError, OSError):
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 14)
        except (IOError, OSError):
            font = ImageFont.load_default()

    for idx, slide in enumerate(slides):
        row = idx // cols
        col = idx % cols

        x = padding + col * (thumb_width + padding)
        y = padding + row * (thumb_height + label_height + padding)

        # Create a placeholder thumbnail (actual rendering needs soffice)
        thumb = _create_slide_placeholder(slide, thumb_width, thumb_height)

        grid.paste(thumb, (x, y))

        # Add label
        label = f"slide{idx + 1}.xml"
        draw.text((x + 5, y + thumb_height + 3), label, fill="black", font=font)

    output_path = f"{output_prefix}.jpg"
    grid.save(output_path, "JPEG", quality=90)
    print(f"Created thumbnail grid: {output_path} ({len(slides)} slides, {cols} cols)")

    return output_path


def _create_slide_placeholder(slide, width: int, height: int) -> Image.Image:
    """Create a placeholder image for a slide showing text content."""
    img = Image.new("RGB", (width, height), "#f0f0f0")
    draw = ImageDraw.Draw(img)

    # Draw border
    draw.rectangle([0, 0, width - 1, height - 1], outline="#cccccc", width=1)

    # Extract text from shapes
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)

    # Draw text preview
    try:
        font = ImageFont.truetype("arial.ttf", 10)
    except (IOError, OSError):
        font = ImageFont.load_default()

    y_pos = 15
    for text in texts[:8]:  # Limit to 8 lines
        truncated = text[:50] + "..." if len(text) > 50 else text
        draw.text((10, y_pos), truncated, fill="#333333", font=font)
        y_pos += 16
        if y_pos > height - 20:
            break

    if not texts:
        draw.text((width // 2 - 30, height // 2), "(empty)", fill="#999999", font=font)

    return img


def _create_thumbnails_fallback(pptx_path, output_prefix, cols, max_slides, thumb_width):
    """Fallback when python-pptx is not available."""
    print("Warning: python-pptx not available. Use soffice + pdftoppm for full rendering:")
    print(f"  python scripts/office/soffice.py --headless --convert-to pdf {pptx_path}")
    print(f"  pdftoppm -jpeg -r 150 output.pdf slide")
    return None


def main():
    parser = argparse.ArgumentParser(description="Create slide thumbnail grid")
    parser.add_argument("pptx_file", help="PPTX file to create thumbnails from")
    parser.add_argument("output_prefix", nargs="?", default="thumbnails", help="Output prefix (default: thumbnails)")
    parser.add_argument("--cols", type=int, default=3, help="Number of columns (default: 3)")
    parser.add_argument("--max-slides", type=int, default=12, help="Max slides per grid (default: 12)")
    args = parser.parse_args()

    create_thumbnails(
        args.pptx_file,
        output_prefix=args.output_prefix,
        cols=args.cols,
        max_slides=args.max_slides,
    )


if __name__ == "__main__":
    main()
