"""
PPT Agent - Utilities

Helper functions for PowerPoint processing based on anthropics/skills PPTX patterns.
Uses python-pptx for all PPTX read/write operations.
"""

import base64
import logging
import os
import time
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from enum import Enum

logger = logging.getLogger("agents.ppt_agent.utils")


# ============================================================================
# COLOR PALETTES (from anthropics/skills PPTX design guidelines)
# ============================================================================

COLOR_PALETTES = {
    "midnight_executive": {
        "primary": "1E2761",
        "secondary": "CADCFC",
        "accent": "FFFFFF",
        "name": "Midnight Executive",
    },
    "forest_moss": {
        "primary": "2C5F2D",
        "secondary": "97BC62",
        "accent": "F5F5F5",
        "name": "Forest & Moss",
    },
    "coral_energy": {
        "primary": "F96167",
        "secondary": "F9E795",
        "accent": "2F3C7E",
        "name": "Coral Energy",
    },
    "warm_terracotta": {
        "primary": "B85042",
        "secondary": "E7E8D1",
        "accent": "A7BEAE",
        "name": "Warm Terracotta",
    },
    "ocean_gradient": {
        "primary": "065A82",
        "secondary": "1C7293",
        "accent": "21295C",
        "name": "Ocean Gradient",
    },
    "charcoal_minimal": {
        "primary": "36454F",
        "secondary": "F2F2F2",
        "accent": "212121",
        "name": "Charcoal Minimal",
    },
    "teal_trust": {
        "primary": "028090",
        "secondary": "00A896",
        "accent": "02C39A",
        "name": "Teal Trust",
    },
    "berry_cream": {
        "primary": "6D2E46",
        "secondary": "A26769",
        "accent": "ECE2D0",
        "name": "Berry & Cream",
    },
    "sage_calm": {
        "primary": "84B59F",
        "secondary": "69A297",
        "accent": "50808E",
        "name": "Sage Calm",
    },
    "cherry_bold": {
        "primary": "990011",
        "secondary": "FCF6F5",
        "accent": "2F3C7E",
        "name": "Cherry Bold",
    },
}

# Typography pairings from anthropics/skills
FONT_PAIRINGS = [
    {"header": "Georgia", "body": "Calibri"},
    {"header": "Arial Black", "body": "Arial"},
    {"header": "Calibri", "body": "Calibri Light"},
    {"header": "Cambria", "body": "Calibri"},
    {"header": "Trebuchet MS", "body": "Calibri"},
    {"header": "Palatino", "body": "Garamond"},
]


# ============================================================================
# READING / TEXT EXTRACTION
# ============================================================================

def extract_text_from_presentation(file_path: str) -> str:
    """Extract all text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text_parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip(" |"):
                        slide_text_parts.append(row_text)

        if slide_text_parts:
            texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text_parts))

    return "\n\n".join(texts)


def analyze_presentation_structure(file_path: str) -> Dict[str, Any]:
    """Analyze the structure and content of a PPTX file."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu

    prs = Presentation(file_path)
    slides_info = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_num,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "shapes": [],
            "text_preview": "",
            "has_notes": False,
        }

        text_parts = []
        for shape in slide.shapes:
            shape_info = {
                "type": shape.shape_type.__class__.__name__ if hasattr(shape, 'shape_type') else "Unknown",
                "name": shape.name,
                "has_text": shape.has_text_frame,
                "has_table": shape.has_table,
            }
            slide_info["shapes"].append(shape_info)

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        text_parts.append(text)

        slide_info["text_preview"] = " | ".join(text_parts[:3])

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes and notes.text.strip():
                slide_info["has_notes"] = True
                slide_info["notes_preview"] = notes.text[:200]

        slides_info.append(slide_info)

    # Presentation-level info
    width = prs.slide_width
    height = prs.slide_height

    return {
        "slide_count": len(prs.slides),
        "slide_width_inches": width / Emu(914400) if width else None,
        "slide_height_inches": height / Emu(914400) if height else None,
        "layout_count": len(prs.slide_layouts),
        "slides": slides_info,
        "file_name": Path(file_path).name,
        "file_size_kb": Path(file_path).stat().st_size / 1024,
    }


# ============================================================================
# CREATION
# ============================================================================

def create_presentation(
    slides_data: List[Dict[str, Any]],
    file_path: str,
    palette_name: str = "midnight_executive",
    font_pairing_index: int = 0,
) -> str:
    """
    Create a new PPTX presentation with design-aware layouts.

    Each slide_data dict can contain:
    - title (str): Slide title
    - content (str or list[str]): Body text or bullet points
    - layout (str): 'title', 'content', 'section', 'two_column', 'blank'
    - notes (str): Speaker notes
    - subtitle (str): For title slides
    - left_content (str or list[str]): For two-column layout
    - right_content (str or list[str]): For two-column layout
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    palette = COLOR_PALETTES.get(palette_name, COLOR_PALETTES["midnight_executive"])
    fonts = FONT_PAIRINGS[font_pairing_index % len(FONT_PAIRINGS)]

    primary_color = RGBColor.from_string(palette["primary"])
    secondary_color = RGBColor.from_string(palette["secondary"])
    accent_color = RGBColor.from_string(palette["accent"])

    for slide_data in slides_data:
        layout_type = slide_data.get("layout", "content")
        title_text = slide_data.get("title", "")
        content = slide_data.get("content", "")
        subtitle = slide_data.get("subtitle", "")
        notes_text = slide_data.get("notes", "")

        if layout_type == "title":
            slide_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(slide_layout)

            if slide.placeholders:
                # Title
                if 0 in slide.placeholders:
                    title_ph = slide.placeholders[0]
                    title_ph.text = title_text
                    for para in title_ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = fonts["header"]
                            run.font.size = Pt(40)
                            run.font.bold = True
                            run.font.color.rgb = accent_color

                # Subtitle
                if 1 in slide.placeholders and subtitle:
                    sub_ph = slide.placeholders[1]
                    sub_ph.text = subtitle
                    for para in sub_ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = fonts["body"]
                            run.font.size = Pt(20)
                            run.font.color.rgb = secondary_color

            # Dark background for title slides
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = primary_color

        elif layout_type == "section":
            slide_layout = prs.slide_layouts[2]  # Section Header
            slide = prs.slides.add_slide(slide_layout)

            if 0 in slide.placeholders:
                title_ph = slide.placeholders[0]
                title_ph.text = title_text
                for para in title_ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = fonts["header"]
                        run.font.size = Pt(36)
                        run.font.bold = True
                        run.font.color.rgb = accent_color

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = primary_color

        elif layout_type == "two_column":
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Title
            from pptx.util import Inches, Pt
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(1)
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = fonts["header"]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = primary_color

            # Left column
            left_content = slide_data.get("left_content", "")
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(4.3), Inches(5)
            )
            _add_content_to_textbox(left_box, left_content, fonts["body"], secondary_color)

            # Right column
            right_content = slide_data.get("right_content", "")
            right_box = slide.shapes.add_textbox(
                Inches(5.2), Inches(1.5), Inches(4.3), Inches(5)
            )
            _add_content_to_textbox(right_box, right_content, fonts["body"], secondary_color)

        elif layout_type == "blank":
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

        else:  # "content" — default
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)

            if 0 in slide.placeholders:
                title_ph = slide.placeholders[0]
                title_ph.text = title_text
                for para in title_ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = fonts["header"]
                        run.font.size = Pt(28)
                        run.font.bold = True
                        run.font.color.rgb = primary_color

            if 1 in slide.placeholders:
                content_ph = slide.placeholders[1]
                if isinstance(content, list):
                    content_ph.text = ""
                    for idx, bullet in enumerate(content):
                        if idx == 0:
                            content_ph.text = bullet
                        else:
                            p = content_ph.text_frame.add_paragraph()
                            p.text = bullet
                        for para in content_ph.text_frame.paragraphs:
                            for run in para.runs:
                                run.font.name = fonts["body"]
                                run.font.size = Pt(16)
                else:
                    content_ph.text = content
                    for para in content_ph.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.name = fonts["body"]
                            run.font.size = Pt(16)

        # Add speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    os.makedirs(os.path.dirname(file_path) or ".", exist_ok=True)
    prs.save(file_path)
    logger.info(f"Created presentation: {file_path} ({len(slides_data)} slides)")
    return file_path


def _add_content_to_textbox(textbox, content, font_name, color):
    """Helper to add content (str or list) to a textbox."""
    from pptx.util import Pt

    tf = textbox.text_frame
    tf.word_wrap = True

    if isinstance(content, list):
        for idx, item in enumerate(content):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.name = font_name
            p.font.size = Pt(14)
    elif content:
        tf.paragraphs[0].text = content
        tf.paragraphs[0].font.name = font_name
        tf.paragraphs[0].font.size = Pt(14)


# ============================================================================
# EDITING
# ============================================================================

def edit_slide_content(
    file_path: str,
    slide_index: int,
    updates: Dict[str, Any],
    output_path: Optional[str] = None,
) -> str:
    """
    Edit content of a specific slide.

    Updates dict can contain:
    - title (str): New title text
    - content (str or list[str]): New body content
    - notes (str): New speaker notes
    """
    from pptx import Presentation

    prs = Presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (0-{len(prs.slides) - 1})")

    slide = prs.slides[slide_index]

    # Update title
    if "title" in updates and 0 in slide.placeholders:
        slide.placeholders[0].text = updates["title"]

    # Update content
    if "content" in updates and 1 in slide.placeholders:
        content = updates["content"]
        content_ph = slide.placeholders[1]
        if isinstance(content, list):
            content_ph.text = ""
            for idx, item in enumerate(content):
                if idx == 0:
                    content_ph.text = item
                else:
                    p = content_ph.text_frame.add_paragraph()
                    p.text = item
        else:
            content_ph.text = content

    # Update notes
    if "notes" in updates:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = updates["notes"]

    save_path = output_path or file_path
    os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
    prs.save(save_path)
    logger.info(f"Updated slide {slide_index + 1} in {save_path}")
    return save_path


def add_slide(
    file_path: str,
    slide_data: Dict[str, Any],
    position: Optional[int] = None,
    output_path: Optional[str] = None,
) -> str:
    """Add a new slide to an existing presentation."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(file_path)

    layout_name = slide_data.get("layout", "content")
    layout_map = {
        "title": 0,
        "content": 1,
        "section": 2,
        "blank": 6,
    }
    layout_idx = layout_map.get(layout_name, 1)

    slide_layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_text = slide_data.get("title", "")
    if title_text and 0 in slide.placeholders:
        slide.placeholders[0].text = title_text

    # Set content
    content = slide_data.get("content", "")
    if content and 1 in slide.placeholders:
        ph = slide.placeholders[1]
        if isinstance(content, list):
            ph.text = ""
            for idx, item in enumerate(content):
                if idx == 0:
                    ph.text = item
                else:
                    p = ph.text_frame.add_paragraph()
                    p.text = item
        else:
            ph.text = content

    # Move to position if specified
    if position is not None and position < len(prs.slides) - 1:
        # python-pptx doesn't natively support reordering, but we can
        # manipulate the XML to move the slide
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        new_slide = slides[-1]
        xml_slides.remove(new_slide)
        xml_slides.insert(position, new_slide)

    save_path = output_path or file_path
    prs.save(save_path)
    logger.info(f"Added slide to {save_path}")
    return save_path


def remove_slide(
    file_path: str,
    slide_index: int,
    output_path: Optional[str] = None,
) -> str:
    """Remove a slide from a presentation by index."""
    from pptx import Presentation
    import copy

    prs = Presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    # Remove slide by manipulating XML
    rId = prs.slides._sldIdLst[slide_index].get("r:id") if hasattr(prs.slides._sldIdLst[slide_index], 'get') else None
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_index])

    save_path = output_path or file_path
    prs.save(save_path)
    logger.info(f"Removed slide {slide_index + 1} from {save_path}")
    return save_path


# ============================================================================
# CONVERSION
# ============================================================================

def convert_slides_to_images(
    file_path: str, output_dir: str, dpi: int = 150
) -> List[str]:
    """
    Convert PPTX slides to images.
    First converts to PDF via LibreOffice, then PDF pages to images.
    Falls back to python-pptx thumbnail extraction if LibreOffice unavailable.
    """
    import subprocess

    os.makedirs(output_dir, exist_ok=True)
    stem = Path(file_path).stem
    image_files = []

    try:
        # Try LibreOffice conversion: PPTX → PDF → Images
        pdf_path = os.path.join(output_dir, f"{stem}.pdf")

        # Convert PPTX to PDF
        subprocess.run(
            [
                "libreoffice", "--headless", "--convert-to", "pdf",
                "--outdir", output_dir, file_path,
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )

        # Convert PDF pages to images
        if os.path.exists(pdf_path):
            try:
                from pdf2image import convert_from_path

                images = convert_from_path(pdf_path, dpi=dpi)
                for i, img in enumerate(images):
                    img_path = os.path.join(output_dir, f"{stem}_slide_{i + 1}.png")
                    img.save(img_path, "PNG")
                    image_files.append(img_path)
            except ImportError:
                # Fall back to pdftoppm if pdf2image not available
                subprocess.run(
                    ["pdftoppm", "-png", "-r", str(dpi), pdf_path,
                     os.path.join(output_dir, f"{stem}_slide")],
                    check=True,
                    capture_output=True,
                )
                for f in sorted(Path(output_dir).glob(f"{stem}_slide*.png")):
                    image_files.append(str(f))

    except (FileNotFoundError, subprocess.CalledProcessError) as e:
        logger.warning(f"LibreOffice conversion failed: {e}")
        # Minimal fallback: create a text representation per slide
        from pptx import Presentation

        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_parts.append(shape.text_frame.text)
            # Save as txt file as minimal fallback
            txt_path = os.path.join(output_dir, f"{stem}_slide_{i + 1}.txt")
            with open(txt_path, "w", encoding="utf-8") as f:
                f.write("\n".join(text_parts))
            image_files.append(txt_path)

    logger.info(f"Converted {len(image_files)} slides to images/text")
    return image_files


# ============================================================================
# CANVAS DISPLAY HELPERS
# ============================================================================

def get_file_base64(file_path: str) -> str:
    """Convert file to base64 for display."""
    with open(file_path, "rb") as f:
        return base64.b64encode(f.read()).decode("utf-8")


def create_pptx_canvas_display(
    file_path: str, title: str
) -> Dict[str, Any]:
    """Create PPTX canvas display for frontend."""
    pptx_base64 = get_file_base64(file_path)
    return {
        "canvas_type": "pptx",
        "title": title,
        "pptx_data": f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{pptx_base64}",
        "file_path": file_path,
        "timestamp": int(time.time() * 1000),
        "no_cache": True,
    }


def get_available_palettes() -> Dict[str, Dict[str, str]]:
    """Return available color palettes for presentation creation."""
    return COLOR_PALETTES


def get_font_pairings() -> List[Dict[str, str]]:
    """Return available font pairings."""
    return FONT_PAIRINGS
